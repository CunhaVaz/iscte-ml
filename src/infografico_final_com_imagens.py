from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Criar uma nova apresentação
prs = Presentation()

# Definir o tamanho do slide como 16:9 widescreen
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Adicionar um slide
layout = prs.slide_layouts[6]  # Layout em branco
slide = prs.slides.add_slide(layout)

# Definir cores
COR_AZUL_ESCURO = RGBColor(0, 32, 96)
COR_AZUL_CLARO = RGBColor(41, 180, 185)
COR_AMARELO = RGBColor(255, 192, 0)
COR_CORAL = RGBColor(255, 111, 97)
COR_CINZA_CLARO = RGBColor(240, 240, 240)
COR_BRANCO = RGBColor(255, 255, 255)

# Adicionar fundo
background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
background.fill.solid()
background.fill.fore_color.rgb = COR_BRANCO
background.line.fill.background()

# Adicionar título principal
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1.0))
title_frame = title_box.text_frame
title_frame.word_wrap = True

p = title_frame.add_paragraph()
p.text = "PREVISÃO DE VENDAS E ANÁLISE DE DRIVERS"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = COR_AZUL_ESCURO
p.alignment = PP_ALIGN.LEFT

p = title_frame.add_paragraph()
p.text = "Análise de Dados para Decisões Estratégicas"
p.font.size = Pt(24)
p.font.color.rgb = COR_AZUL_ESCURO
p.alignment = PP_ALIGN.LEFT

# Adicionar linha divisória horizontal
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.7), Inches(12.33), Inches(0.03))
line.fill.solid()
line.fill.fore_color.rgb = COR_AZUL_CLARO
line.line.fill.background()

# Seção 1: Objetivo
left = Inches(0.5)
top = Inches(2.0)
width = Inches(3.0)
height = Inches(2.0)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = COR_CINZA_CLARO
shape.line.color.rgb = COR_AZUL_CLARO
shape.line.width = Pt(1)

# Adicionar ícone para Objetivo
icon_path = "/home/ubuntu/infografico_images/data_science_icon.png"
if os.path.exists(icon_path):
    icon_left = left + Inches(1.0)
    icon_top = top + Inches(0.3)
    icon_width = Inches(1.0)
    icon_height = Inches(1.0)
    slide.shapes.add_picture(icon_path, icon_left, icon_top, icon_width, icon_height)

# Título da seção
title_box = slide.shapes.add_textbox(left, top + Inches(0.1), width, Inches(0.4))
title_frame = title_box.text_frame
p = title_frame.add_paragraph()
p.text = "OBJETIVO"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COR_AZUL_ESCURO
p.alignment = PP_ALIGN.CENTER

# Conteúdo da seção
content_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.4), width - Inches(0.4), Inches(0.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

objectives = [
    "• Prever vendas por período/cliente",
    "• Identificar principais drivers",
    "• Otimizar estratégias com dados"
]

for obj in objectives:
    p = content_frame.add_paragraph()
    p.text = obj
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.LEFT

# Seção 2: Pipeline
left = Inches(3.8)
top = Inches(2.0)
width = Inches(3.0)
height = Inches(2.0)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = COR_CINZA_CLARO
shape.line.color.rgb = COR_AZUL_CLARO
shape.line.width = Pt(1)

# Adicionar ícone para Pipeline
icon_path = "/home/ubuntu/infografico_images/analytics_icon.png"
if os.path.exists(icon_path):
    icon_left = left + Inches(1.0)
    icon_top = top + Inches(0.3)
    icon_width = Inches(1.0)
    icon_height = Inches(1.0)
    slide.shapes.add_picture(icon_path, icon_left, icon_top, icon_width, icon_height)

# Título da seção
title_box = slide.shapes.add_textbox(left, top + Inches(0.1), width, Inches(0.4))
title_frame = title_box.text_frame
p = title_frame.add_paragraph()
p.text = "PIPELINE"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COR_AZUL_ESCURO
p.alignment = PP_ALIGN.CENTER

# Conteúdo da seção
content_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.4), width - Inches(0.4), Inches(0.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

steps = [
    "1. Limpeza → dados omissos",
    "2. EDA → estatísticas e gráficos",
    "3. Modelos → Linear vs RF",
    "4. Métricas → RMSE, MAPE, R²",
    "5. Insights → drivers + previsões"
]

for step in steps:
    p = content_frame.add_paragraph()
    p.text = step
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.LEFT

# Seção 3: Modelo Vencedor
left = Inches(7.1)
top = Inches(2.0)
width = Inches(3.0)
height = Inches(2.0)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = COR_CINZA_CLARO
shape.line.color.rgb = COR_AZUL_CLARO
shape.line.width = Pt(1)

# Adicionar gráfico para Modelo Vencedor
chart_path = "/home/ubuntu/upload/search_images/sEjjt7arToeP.png"
if os.path.exists(chart_path):
    chart_left = left + Inches(0.5)
    chart_top = top + Inches(0.3)
    chart_width = Inches(2.0)
    chart_height = Inches(1.0)
    slide.shapes.add_picture(chart_path, chart_left, chart_top, chart_width, chart_height)

# Título da seção com fundo amarelo
title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.5), top + Inches(0.1), width - Inches(1.0), Inches(0.4))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = COR_AMARELO
title_shape.line.color.rgb = COR_AMARELO
title_shape.shadow.inherit = False

title_box = slide.shapes.add_textbox(left + Inches(0.5), top + Inches(0.1), width - Inches(1.0), Inches(0.4))
title_frame = title_box.text_frame
p = title_frame.add_paragraph()
p.text = "MODELO VENCEDOR"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COR_AZUL_ESCURO
p.alignment = PP_ALIGN.CENTER

# Conteúdo da seção
content_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.4), width - Inches(0.4), Inches(0.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

p = content_frame.add_paragraph()
p.text = "Random Forest"
p.font.size = Pt(14)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

metrics = [
    "RMSE (teste) ≈ 3000",
    "MAPE (teste) ≈ 7%",
    "R² (teste) ≈ 0,93"
]

for metric in metrics:
    p = content_frame.add_paragraph()
    p.text = metric
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.LEFT

# Seção 4: Variáveis-chave
left = Inches(10.4)
top = Inches(2.0)
width = Inches(2.5)
height = Inches(2.0)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = COR_CINZA_CLARO
shape.line.color.rgb = COR_AZUL_CLARO
shape.line.width = Pt(1)

# Adicionar ícone para Variáveis-chave
icon_path = "/home/ubuntu/infografico_images/business_analyst.jpg"
if os.path.exists(icon_path):
    icon_left = left + Inches(0.75)
    icon_top = top + Inches(0.3)
    icon_width = Inches(1.0)
    icon_height = Inches(1.0)
    slide.shapes.add_picture(icon_path, icon_left, icon_top, icon_width, icon_height)

# Título da seção
title_box = slide.shapes.add_textbox(left, top + Inches(0.1), width, Inches(0.4))
title_frame = title_box.text_frame
p = title_frame.add_paragraph()
p.text = "VARIÁVEIS-CHAVE"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COR_AZUL_ESCURO
p.alignment = PP_ALIGN.CENTER

# Conteúdo da seção
content_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.4), width - Inches(0.4), Inches(0.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

variables = [
    "1. Margem_Valor",
    "2. Produto",
    "3. Cliente",
    "*Regra 80/20 confirmada*"
]

for var in variables:
    p = content_frame.add_paragraph()
    p.text = var
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.LEFT

# Seção 5: Valor para Decisão
left = Inches(0.5)
top = Inches(4.3)
width = Inches(6.3)
height = Inches(2.0)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = COR_CINZA_CLARO
shape.line.color.rgb = COR_AZUL_CLARO
shape.line.width = Pt(1)

# Adicionar gráfico para Valor para Decisão
chart_path = "/home/ubuntu/infografico_images/sales_forecast.png"
if os.path.exists(chart_path):
    chart_left = left + Inches(0.5)
    chart_top = top + Inches(0.5)
    chart_width = Inches(2.5)
    chart_height = Inches(1.4)
    slide.shapes.add_picture(chart_path, chart_left, chart_top, chart_width, chart_height)

# Título da seção
title_box = slide.shapes.add_textbox(left, top + Inches(0.1), width, Inches(0.4))
title_frame = title_box.text_frame
p = title_frame.add_paragraph()
p.text = "VALOR PARA DECISÃO"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COR_AZUL_ESCURO
p.alignment = PP_ALIGN.CENTER

# Conteúdo da seção
content_box = slide.shapes.add_textbox(left + Inches(3.3), top + Inches(0.5), width - Inches(3.5), Inches(1.4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

values = [
    "• Previsões fiáveis (~7% erro)",
    "• Foco em clientes/produtos-chave",
    "• Apoio ao planeamento",
    "• Precisão do modelo: 93%"
]

for value in values:
    p = content_frame.add_paragraph()
    p.text = value
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.LEFT

# Seção 6: Reflexão
left = Inches(7.1)
top = Inches(4.3)
width = Inches(5.8)
height = Inches(2.0)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = COR_CINZA_CLARO
shape.line.color.rgb = COR_AZUL_CLARO
shape.line.width = Pt(1)

# Adicionar ícone para Reflexão
icon_path = "/home/ubuntu/infografico_images/data_science_icon.png"
if os.path.exists(icon_path):
    icon_left = left + Inches(0.5)
    icon_top = top + Inches(0.5)
    icon_width = Inches(1.0)
    icon_height = Inches(1.0)
    slide.shapes.add_picture(icon_path, icon_left, icon_top, icon_width, icon_height)

# Título da seção
title_box = slide.shapes.add_textbox(left, top + Inches(0.1), width, Inches(0.4))
title_frame = title_box.text_frame
p = title_frame.add_paragraph()
p.text = "REFLEXÃO E PRÓXIMOS PASSOS"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COR_AZUL_ESCURO
p.alignment = PP_ALIGN.CENTER

# Conteúdo da seção em duas colunas
# Coluna 1
col1_box = slide.shapes.add_textbox(left + Inches(1.7), top + Inches(0.5), width/2 - Inches(1.7), Inches(1.4))
col1_frame = col1_box.text_frame
col1_frame.word_wrap = True

p = col1_frame.add_paragraph()
p.text = "Obstáculos:"
p.font.size = Pt(12)
p.font.bold = True
p.alignment = PP_ALIGN.LEFT

obstacles = [
    "• Qualidade dos dados",
    "• Concentração de vendas",
    "• Sazonalidade"
]

for obs in obstacles:
    p = col1_frame.add_paragraph()
    p.text = obs
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.LEFT

# Coluna 2
col2_box = slide.shapes.add_textbox(left + width/2, top + Inches(0.5), width/2 - Inches(0.3), Inches(1.4))
col2_frame = col2_box.text_frame
col2_frame.word_wrap = True

p = col2_frame.add_paragraph()
p.text = "Próximos Passos:"
p.font.size = Pt(12)
p.font.bold = True
p.alignment = PP_ALIGN.LEFT

next_steps = [
    "• ARIMA/Prophet para séries temporais",
    "• Dashboards executivos interativos",
    "• Automatização do processo"
]

for step in next_steps:
    p = col2_frame.add_paragraph()
    p.text = step
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.LEFT

# Adicionar elementos decorativos
# Canto inferior direito
corner_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.33), Inches(6.5), Inches(2), Inches(1))
corner_shape.fill.solid()
corner_shape.fill.fore_color.rgb = COR_AZUL_CLARO
corner_shape.line.fill.background()
corner_shape.rotation = 0

# Salvar a apresentação
prs.save('/home/ubuntu/infografico_final_com_imagens.pptx')